# -*- coding: utf-8 -*-
"""
장애 주간보고 - PPT 생성 프로그램
=================================

update_excel.py 로 갱신된 엑셀을 읽어, 지난주 장애공유 PPT(=템플릿 겸 비교본)의
차트/요약문구/제목을 채워 이번주 주간보고 PPT를 만든다.

갱신 대상
  - 슬라이드 0: 전사 장애 발생 추이 차트 + 요약문구(요약!C22/C23) + 증감문구(요약!C24)
  - 슬라이드 1: 계열사 전사 추이 차트 + 요약문구(요약!C135/C136) + 증감문구(요약!C137)
  - 슬라이드 2: 사업부 4개 추이 차트 (사업부별 시트)
  - 슬라이드 3: 사업부 4개 계열사 추이 차트 (사업부별_계열사별 시트)
  - 슬라이드 6~9: 제목의 등급 건수 (요약!C22에서 총/2등급/3등급이하 추출)
  - 슬라이드 4 '신규 등급 확정 장애' 표: 지난주 PPT(--prev)와 비교해 자동 작성
      · 이번주 엑셀 '데이터'(등급확정) 중 지난주 상세목록(6p~)에 없는 건 = 신규 확정
      · 장애명 = [등급]+제목
      · 장애내용 = 지난주 PPT(4p 신규발생 → 없으면 5p 협의중)에서 찾은 값과 엑셀 N열 값 중 '더 긴' 텍스트
      · 나머지 칸(발생일/시간/고객/원인/조직)은 엑셀 '데이터'에서 채움
      · 동일 장애 매칭 = 발생일(월/일) + (고객사 부분일치 OR 제목 핵심 일치)
  - 슬라이드 4 '신규 발생 장애' / 5 '등급 협의 중 장애' 표: 이번주 등급확정으로 넘어간 건은 삭제
  - 슬라이드 6~ '장애 목록 상세': 4p 신규확정 중 목록에 없는 건 추가 후 (등급→발생일) 재정렬
      · 정렬 = 등급 오름차순(1등급 최상단) → 같은 등급 내 발생일 오름차순
      · 한 페이지 용량 초과 시 상세 페이지를 자동 복제·추가
      · 제목 뒤 등급건수 텍스트는 14pt
  - 모든 표: 내부 행 구분선을 가늘고 연하게 통일(행 추가/삭제로 생긴 두꺼운 선 정리, 바깥 테두리는 유지)
  - 각 슬라이드의 '(M/D 기준)' 등 날짜 표기를 기준일로 갱신
  * 차트는 작성일 기준 최근 3개년(= 기준연도-2년 1월 ~ 기준월) 을 그림

요구사항: python-pptx, openpyxl
사용법 (--prev 필수):
    python make_ppt.py --prev "지난주_장애공유.pptx"
    python make_ppt.py --date 2026-08-05 --prev "지난주_장애공유.pptx"
"""

import argparse
import copy
import datetime as dt
import glob
import os
import re
import sys

import openpyxl
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

import appconfig   # 회사 고유값(사업부명 등)은 report_config.json 에서 로드

# ─────────────────────────────────────────────────────────────────────────────
# 엑셀 시트 레이아웃 (update_excel.py 와 동일 규칙: 2024년 1월 시작)
# ─────────────────────────────────────────────────────────────────────────────
BASE_YEAR = 2024

# (시트명, 전사 데이터행, 전사 시작열, 사업부별 시작열, {사업부: 데이터행})
SRC = {
    "사업부별": {
        "sheet": "사업부별",
        "total_row": 23, "total_start_col": 2,     # 전사: B(2024-01)
        "div_start_col": 41,                         # AO(2024-01)
        "divisions": appconfig.DIVISIONS_MAIN,       # {사업부명: 데이터행} (config)
    },
    "사업부별_계열사별": {
        "sheet": "사업부별_계열사별",
        "total_row": 17, "total_start_col": 1,       # 전사(계열사): A(2024-01)
        "div_start_col": 39,                          # AM(2024-01)
        "divisions": appconfig.DIVISIONS_AFFILIATE,  # {사업부명: 데이터행} (config)
    },
}


def log(m):
    print(m, flush=True)


def col_for_month(start_col, year, month):
    return start_col + (year - BASE_YEAR) * 12 + (month - 1)


def month_range(y0, m0, y1, m1):
    out = []
    y, m = y0, m0
    while (y, m) <= (y1, m1):
        out.append((y, m))
        m += 1
        if m > 12:
            m, y = 1, y + 1
    return out


def read_series(ws, row, start_col, months):
    vals = []
    for (y, m) in months:
        c = col_for_month(start_col, y, m)
        v = ws.cell(row=row, column=c).value
        vals.append(0 if v is None else v)
    return vals


def cell(ws, coord):
    v = ws[coord].value
    return "" if v is None else str(v)


def group_commas(text):
    """4자리 이상 정수에 천단위 콤마 삽입. 단, '~년'(연도) 표기는 제외."""
    return re.sub(r"(?<!\d)(\d{4,})(?!\d)(?!\s*년)",
                  lambda m: format(int(m.group(1)), ","), text)


def polish_delta(text):
    """증감 문구 정리:
       - 백분율 앞의 부호(-) 제거
       - 방향 표기를 부호로 결정 (음수 -> '감소', 그 외 -> '증가'). '감소/증가' 모호 표기도 정리
       - 마지막에 천단위 콤마 적용
    """
    def repl(m):
        sign, num = m.group(1), m.group(2)
        direction = "감소" if sign == "-" else "증가"
        return f"{num}% {direction}"

    text = re.sub(r"([+-]?)(\d+(?:\.\d+)?)\s*%\s*(?:감소/증가|감소|증가)", repl, text)
    return group_commas(text)


# ─────────────────────────────────────────────────────────────────────────────
# 텍스트 갱신 헬퍼 (첫 run 서식을 유지하며 텍스트 교체)
# ─────────────────────────────────────────────────────────────────────────────
def set_paragraph_text(p, text):
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        run = p.add_run()
        run.text = text


def set_shape_text(shape, text):
    """단일 문단 텍스트박스의 텍스트를 첫 run 서식 유지하며 교체"""
    set_paragraph_text(shape.text_frame.paragraphs[0], text)


def sub_date(text, base):
    """텍스트 내 'M/D' 및 '’YY' 날짜 표기를 기준일로 치환"""
    text = re.sub(r"\d{1,2}/\d{1,2}", f"{base.month}/{base.day}", text)
    text = re.sub(r"’\d\d", f"’{base:%y}", text)
    return text


def find_shape(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    return None


def _num(v):
    try:
        f = float(v)
        return str(int(f)) if f.is_integer() else str(f)
    except Exception:
        return str(v)


def _col_idx(col):
    n = 0
    for ch in col:
        n = n * 26 + (ord(ch) - 64)
    return n


def _col_letter(n):
    s = ""
    while n > 0:
        n, r = divmod(n - 1, 26)
        s = chr(65 + r) + s
    return s


def _extend_f(f_el, npts):
    """외부 참조식 c:f 의 끝 열을 포인트 수에 맞게 확장 (표시는 캐시 사용, 참조 정합성용)"""
    if f_el is None or not f_el.text:
        return
    m = re.match(r"^(.*!)\$?([A-Z]+)\$?(\d+):\$?([A-Z]+)\$?(\d+)$", f_el.text)
    if not m:
        return
    sheet, c1, r1, c2, r2 = m.groups()
    end = _col_letter(_col_idx(c1) + npts - 1)
    f_el.text = f"{sheet}${c1}${r1}:${end}${r2}"


def _rebuild_cache(ref_el, cache_tag, items, is_num):
    """c:strCache / c:numCache 의 ptCount + pt 를 새 값으로 재작성"""
    cache = ref_el.find(qn(cache_tag))
    if cache is None:
        return
    fmt = cache.find(qn("c:formatCode")) if is_num else None
    for el in list(cache):
        if el.tag in (qn("c:ptCount"), qn("c:pt")):
            cache.remove(el)
    pc = cache.makeelement(qn("c:ptCount"), {"val": str(len(items))})
    if fmt is not None:
        fmt.addnext(pc)
    else:
        cache.insert(0, pc)
    prev = pc
    for i, item in enumerate(items):
        pt = cache.makeelement(qn("c:pt"), {"idx": str(i)})
        v = cache.makeelement(qn("c:v"), {})
        v.text = _num(item) if is_num else str(item)
        pt.append(v)
        prev.addnext(pt)
        prev = pt


def align_year_bands(slide, n_tpl, n_new, base_year):
    """슬라이드 0/1: 연도 구분 회색 음영(사각형)과 연도 라벨을 1월 경계에 맞춰 재배치.
    템플릿의 기존 밴드 위치로 플롯 기하(좌측/월폭)를 자가 보정한 뒤, 새 포인트 수(n_new)로 재계산."""
    rects = sorted([sh for sh in slide.shapes if sh.shape_type == 1], key=lambda s: int(s.left))
    year_labels = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            m = re.fullmatch(r"(\d{4})년", sh.text_frame.text.strip())
            if m:
                year_labels.append((int(m.group(1)), sh))
    if len(rects) < 2 or n_tpl < 2:
        return
    band_a, band_b = rects[0], rects[1]           # 2번째 연도 밴드, 3번째 연도 밴드
    bw_tpl = (band_b.left - band_a.left) / 12.0    # 템플릿 기준 1개월 폭
    if bw_tpl <= 0:
        return
    inner_left = band_a.left - 12 * bw_tpl          # 플롯 영역 좌측(=2024-01 경계)
    inner_width = n_tpl * bw_tpl                     # 플롯 영역 가로폭(포인트 수와 무관, 물리적 고정)
    bw = inner_width / n_new                         # 새 1개월 폭

    def jan(k):
        return int(round(inner_left + k * bw))

    # 회색 음영: (base_year-1) 연도 = 12개월, base_year 연도 = 나머지 개월
    band_a.left = jan(12)
    band_a.width = int(round(12 * bw))
    band_b.left = jan(24)
    band_b.width = int(round((n_new - 24) * bw))
    # 연도 라벨: 각 연도의 1월 경계에 좌측 정렬
    for year, sh in year_labels:
        k = (year - (base_year - 2)) * 12
        if 0 <= k <= n_new:
            sh.left = jan(k)


def chart_inner_xw(chart):
    """차트 plotArea inner manualLayout 의 (x비율, w비율). 없으면 기본값(다른 차트와 동일)."""
    default = (0.032337962962962964, 0.93532407407407403)
    try:
        cs = chart._chartSpace
        chart_el = cs.find(qn("c:chart"))
        plot_area = chart_el.find(qn("c:plotArea"))
        layout = plot_area.find(qn("c:layout"))
        ml = layout.find(qn("c:manualLayout")) if layout is not None else None
        if ml is None:
            return default
        x = ml.find(qn("c:x"))
        w = ml.find(qn("c:w"))
        return (float(x.get("val")), float(w.get("val")))
    except Exception:
        return default


def align_year_labels_grid(slide, n_new, base_year):
    """슬라이드 2/3: 각 사업부 차트 하단의 연도 라벨을 해당 차트의 1월 경계에 좌측 정렬."""
    charts = [sh for sh in slide.shapes if sh.has_chart]
    moved = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        m = re.fullmatch(r"(\d{4})년", sh.text_frame.text.strip())
        if not m:
            continue
        year = int(m.group(1))
        cx = sh.left + sh.width // 2
        host = None
        for ch in charts:
            if ch.left <= cx <= ch.left + ch.width and ch.top <= sh.top <= ch.top + ch.height:
                host = ch
                break
        if host is None:
            continue
        xf, wf = chart_inner_xw(host.chart)
        inner_left = host.left + xf * host.width
        bw = (wf * host.width) / n_new
        k = (year - (base_year - 2)) * 12
        if 0 <= k <= n_new:
            # 가운데 정렬 + 넓은 박스라 글자가 우측으로 밀리므로, 왼쪽 정렬 + 여백 0 으로 바꿔
            # 박스 왼쪽(=1월 칸 시작 ≈ '1월' 텍스트 시작점)에 글자가 시작되도록 정렬
            tf = sh.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
            sh.left = int(round(inner_left + k * bw))
            moved += 1
    return moved


def replace_chart(shape, categories, values, series_name=""):
    """외부 링크 차트: 시리즈의 카테고리/값 캐시를 직접 재작성 (트렌드라인·서식 보존)"""
    ser = shape.chart.series[0]._element
    cat = ser.find(qn("c:cat"))
    val = ser.find(qn("c:val"))
    if cat is not None:
        strRef = cat.find(qn("c:strRef"))
        if strRef is not None:
            _rebuild_cache(strRef, "c:strCache", categories, is_num=False)
            _extend_f(strRef.find(qn("c:f")), len(categories))
        else:
            numRef = cat.find(qn("c:numRef"))
            if numRef is not None:
                _rebuild_cache(numRef, "c:numCache", categories, is_num=False)
                _extend_f(numRef.find(qn("c:f")), len(categories))
    if val is not None:
        numRef = val.find(qn("c:numRef"))
        if numRef is not None:
            _rebuild_cache(numRef, "c:numCache", values, is_num=True)
            _extend_f(numRef.find(qn("c:f")), len(values))


# ─────────────────────────────────────────────────────────────────────────────
# 슬라이드 4: '신규 등급 확정 장애' 표 작성 (지난주 PPT 비교)
#   엑셀 '데이터'(이번주 등급확정) 중 지난주 상세목록(6p~)에 없는 건 = 신규 확정.
#   장애명 = [등급]+제목,  장애내용 = 지난주 4p 신규발생 → 없으면 5p 협의중 → 없으면 엑셀,
#   나머지 칸(발생일/시간/고객/원인/조직) = 엑셀 데이터.
#   매칭 기준 = 발생일(월/일) + 고객사(부분일치).
# ─────────────────────────────────────────────────────────────────────────────
# '데이터' 시트 컬럼 번호
DATA_COLS = dict(year=1, title=3, grade=4, start=9, mins=11, cust=13,
                 content=14, cause=15, div=17, charge=18, team=19)


# ─────────────────────────────────────────────────────────────────────────────
# 장애 '내용' 요약 (규칙 기반 추출 · AI 미사용)
#   엑셀 '장애 내용'은 대부분 (1) '문자발송 정형양식'(■ 발생일시/고객사명/…/장애내용:…),
#   (2) '1.장애요약 …' 번호 섹션형, (3) 자유 서술 이다.
#   → 핵심 필드('장애내용'/'요약')를 추출하거나, 없으면 앞 문장을 취해 표에 넣기 좋게 축약.
#   시간대별 사건 로그(타임라인)형은 규칙으로 압축이 어려워 '저신뢰'로 표시(수동 검토 유도).
#   ※ 외부 AI 호출 없음 — 사내/고객 정보가 외부로 나가지 않음.
# ─────────────────────────────────────────────────────────────────────────────
_SUMM_LABELS = (r'발생일시|일시|발생일|복구|고객사명|고객사|시스템명|장애\s*발생\s*시스템|장비|'
                r'장애\s*유형|장애유형|비즈영향도|영향도|추정등급|등급|발신자|발송조직|팀장명|'
                r'장애\s*내용|장애내용|내용|조치\s*내용|조치내용|조치사항|조치|장애\s*원인|원인|'
                r'현상|영향|재발\s*방지|비고|특이사항|대상|경위')


def _summ_clean(s):
    s = s.replace('\r', '\n').replace('_x000B_', '\n')
    s = re.sub(r'https?://\S+|www\.\S+', ' ', s)      # URL 제거
    s = re.sub(r'[■▶◆●□▪]', '\n', s)                  # 불릿 → 줄바꿈(필드 경계)
    s = re.sub(r'[ \t]+', ' ', s)
    s = re.sub(r'\n{2,}', '\n', s)
    return s.strip()


def _summ_field(s, name):
    """'name : 값' 필드에서 값을 추출(다음 필드 라벨 직전까지)."""
    m = re.search(r'(?:^|\n)\s*' + name + r'\s*[:：]\s*(.+?)'
                  r'(?=\n\s*(?:' + _SUMM_LABELS + r')\s*[:：]|\Z)', s, re.S)
    return re.sub(r'\s+', ' ', m.group(1)).strip(' -·•*:') if m else None


def _summ_sentences(s):
    s = s.replace('\n', ' ')
    parts = re.split(r'(?<=[다함음됨임])\.\s|(?<=[.!?。])\s', s)
    return [p.strip(' ·-•*:') for p in parts if p.strip(' ·-•*:')]


def _summ_trim(t, n=140):
    t = re.sub(r'\s+', ' ', t).strip(' -·•*:').rstrip(' .,')
    if len(t) <= n:
        return t
    cut = t[:n]
    b = max(cut.rfind('. '), cut.rfind(', '), cut.rfind(' '))
    return (cut[:b] if b > 70 else cut).rstrip(' ,.') + '…'


def summarize_content(raw):
    """엑셀 '장애 내용' → 표에 넣을 요약. 반환: (요약문, 저신뢰여부)."""
    raw = (raw or '').strip()
    if len(raw) < 80:                       # 이미 짧으면 그대로 사용
        return raw, False
    s = _summ_clean(raw)
    # 1) '장애내용' 필드 우선(+짧으면 '조치' 결합)
    content = _summ_field(s, r'장애\s*내용') or _summ_field(s, r'내용')
    if content:
        act = (_summ_field(s, r'조치\s*내용') or _summ_field(s, r'조치사항')
               or _summ_field(s, r'조치'))
        out = content
        if act and len(out) < 90 and act not in out:
            out = f"{out} (조치: {act})"
        return _summ_trim(out), False
    # 2) '장애요약' 섹션
    m = re.search(r'(?:\d+\s*[.)]\s*)?(장애\s*요약|요약)\s*[:\-)]?\s*', s)
    if m:
        start = m.end()
        nxt = re.search(r'\n\s*(?:\d+\s*[.)]\s*)|\n\s*(?:' + _SUMM_LABELS + r')\s*[:：]',
                        s[start:])
        blk = (s[start:start + nxt.start()] if nxt else s[start:])
        blk = re.sub(r'대상\s*도메인.*', '', blk, flags=re.S).strip(' \n:-)')
        if blk:
            return _summ_trim(blk), False
    # 3) 폴백: 앞 1~2문장 (시간대 로그=타임라인형이면 '저신뢰'로 표시)
    ss = _summ_sentences(s)
    out = ss[0] if ss else s
    if ss and len(out) < 30 and len(ss) > 1:
        out = out + ' ' + ss[1]
    low_conf = len(re.findall(r'\d{1,2}:\d{2}', s)) >= 2
    return _summ_trim(out), low_conf


def _norm_date(s):
    """'M/D' 또는 datetime -> (월, 일)"""
    if isinstance(s, dt.datetime):
        return (s.month, s.day)
    m = re.search(r"(\d{1,2})\s*/\s*(\d{1,2})", str(s or ""))
    return (int(m.group(1)), int(m.group(2))) if m else None


def _cust_match(a, b):
    """고객사 부분일치 (공백 제거 후 포함관계)"""
    a = re.sub(r"\s", "", a or "")
    b = re.sub(r"\s", "", b or "")
    return bool(a) and bool(b) and (a in b or b in a)


def _title_core(s):
    """장애명에서 앞쪽 '[...]' 프리픽스와 공백을 제거한 핵심 문자열"""
    s = re.sub(r"^\s*(\[[^\]]*\]\s*)+", "", s or "")
    return re.sub(r"\s", "", s)


def _same_incident(date, cust, title, rec):
    """동일 장애 판정: 발생일 일치 + (고객사 부분일치 OR 제목 핵심 포함관계).
    엑셀과 지난주 PPT는 고객사/제목 표기가 제각각이라 둘 중 하나만 맞아도 동일로 본다."""
    if rec["date"] != date:
        return False
    if _cust_match(cust, rec["cust"]):
        return True
    t1, t2 = _title_core(title), _title_core(rec["name"])
    return len(t1) >= 4 and len(t2) >= 4 and (t1 in t2 or t2 in t1)


def _find_tables_by_title(prs, keyword):
    """제목 텍스트박스에 keyword 가 포함된 슬라이드에서, 그 제목 바로 아래 표들을 반환."""
    out = []
    for slide in prs.slides:
        titles = [sh for sh in slide.shapes
                  if sh.has_text_frame and keyword in sh.text_frame.text]
        tables = [sh for sh in slide.shapes if sh.has_table]
        for tsh in titles:
            below = [tb for tb in tables if tb.top is not None and tb.top >= tsh.top]
            if below:
                out.append(min(below, key=lambda tb: tb.top - tsh.top).table)
    return out


def _table_records(table):
    """표의 데이터행(2행부터) -> [{name,date,cust,content}]"""
    recs = []
    for ri in range(2, len(table.rows)):
        cells = table.rows[ri].cells
        name = cells[0].text.strip()
        if not name:
            continue
        recs.append({"name": name, "date": _norm_date(cells[1].text),
                     "cust": cells[3].text.strip(), "content": cells[4].text.strip()})
    return recs


def _set_table_data_rows(table, n_data, header_rows=2):
    """데이터 행 수를 n_data 로 맞춤(마지막 데이터행 복제/삭제로 서식 유지)."""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    cur = len(trs) - header_rows
    if cur < n_data:
        src = trs[-1]
        for _ in range(n_data - cur):
            tbl.append(copy.deepcopy(src))
    elif cur > n_data:
        for tr in trs[header_rows + n_data:]:
            tbl.remove(tr)


def _est_lines(text, chars_per_line):
    """셀 텍스트의 대략적인 줄 수(개행 + 폭 기준 줄바꿈)"""
    n = 0
    for para in str(text or "").split("\n"):
        n += max(1, -(-len(para) // max(1, chars_per_line)))
    return max(1, n)


# 신규확정/신규발생 표 9개 열의 대략 폭(열당 한 줄에 들어가는 글자 수)
_TABLE_CPL = [10, 5, 5, 6, 26, 6, 7, 7, 6]


def _est_table_height(table):
    """표의 대략 렌더 높이(EMU) 추정 — 각 행은 가장 긴 셀 줄 수 기준."""
    LINE, PAD = 210000, 120000
    total = 0
    for row in table.rows:
        cells = row.cells
        lines = 1
        for ci in range(min(len(_TABLE_CPL), len(cells))):
            lines = max(lines, _est_lines(cells[ci].text, _TABLE_CPL[ci]))
        total += lines * LINE + PAD
    return total


def _fill_txbody(txBody, text):
    """txBody(a:txBody)의 텍스트 교체(첫 run 서식 유지, 개행은 문단 분리)."""
    text = "" if text is None else str(text).replace("\r", "")
    lines = text.split("\n") or [""]
    ps = txBody.findall(qn("a:p"))
    tmpl_p = ps[0]
    tmpl_r = tmpl_p.find(qn("a:r"))
    tmpl_rpr = tmpl_r.find(qn("a:rPr")) if tmpl_r is not None else None
    for p in ps:
        txBody.remove(p)
    for line in lines:
        p = copy.deepcopy(tmpl_p)
        for ch in list(p):
            if ch.tag in (qn("a:r"), qn("a:br"), qn("a:fld"), qn("a:endParaRPr")):
                p.remove(ch)
        r = p.makeelement(qn("a:r"), {})
        if tmpl_rpr is not None:
            r.append(copy.deepcopy(tmpl_rpr))
        t = r.makeelement(qn("a:t"), {})
        t.text = line
        r.append(t)
        p.append(r)
        txBody.append(p)


def _set_cell_text(cell, text):
    _fill_txbody(cell.text_frame._txBody, text)


def _set_tc_text(tc, text):
    tb = tc.find(qn("a:txBody"))
    if tb is not None:
        _fill_txbody(tb, text)


def _grade_of(name):
    m = re.match(r"\s*\[(\d)\s*등급\]", name or "")
    return int(m.group(1)) if m else 9


def _tr_est_height(tr):
    """행(a:tr)의 대략 높이(EMU) 추정 — 가장 높은 셀 기준(줄 수 + 문단/불릿 간격 반영)."""
    best = 0
    for ci, tc in enumerate(tr.findall(qn("a:tc"))):
        if ci >= len(_TABLE_CPL):
            break
        tb = tc.find(qn("a:txBody"))
        if tb is None:
            continue
        lines, paras = 0, 0
        for p in tb.findall(qn("a:p")):
            ptxt = "".join((t.text or "") for t in p.iter(qn("a:t")))
            lines += _est_lines(ptxt, _TABLE_CPL[ci])
            paras += 1
        cell_h = max(1, lines) * 200000 + max(0, paras - 1) * 55000  # 불릿 간 간격 가중
        best = max(best, cell_h)
    return best + 110000


def _detail_slide_tables(prs):
    """'장애 목록 상세' 슬라이드들의 (slide, table_shape)을 슬라이드 순서대로 반환."""
    out = []
    for slide in prs.slides:
        has_title = any(sh.has_text_frame and "장애 목록 상세" in sh.text_frame.text
                        for sh in slide.shapes)
        tables = [sh for sh in slide.shapes if sh.has_table]
        if has_title and tables:
            out.append((slide, min(tables, key=lambda tb: tb.top)))
    return out


def _duplicate_detail_slide(prs, src_slide):
    """상세 슬라이드를 복제해 맨 뒤에 추가하고 (slide, table_shape) 반환.
    상세 슬라이드는 표/텍스트/선만 있어 외부 관계(rel)가 없으므로 도형 XML 복사로 충분."""
    dest = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(dest.shapes):          # 레이아웃 placeholder 제거
        shp._element.getparent().remove(shp._element)
    for shp in src_slide.shapes:           # 원본 도형 복제
        dest.shapes._spTree.append(copy.deepcopy(shp._element))
    tsh = min([sh for sh in dest.shapes if sh.has_table], key=lambda tb: tb.top)
    return dest, tsh


def _clear_data_rows(table):
    tbl = table._tbl
    for tr in tbl.findall(qn("a:tr"))[2:]:
        tbl.remove(tr)


def rebuild_detail_list(prs, new_items, confirmed=None):
    """6p~ '장애 목록 상세'에 신규확정 누락건을 추가하고 (등급 asc, 발생일 asc)로 재정렬·재배치.
    한 페이지 용량을 넘으면 상세 페이지를 복제해 추가한다. 추가된 건수 반환.
    confirmed(이번주 엑셀 등급확정 목록)가 주어지면, 엑셀에서 사라진(등급확정 해제/삭제) 기존 행은
    상세목록에서 제거한다(프루닝)."""
    dinfo = _detail_slide_tables(prs)
    if not dinfo:
        return 0
    entries = []      # [grade, (m,d), tr_element]
    existing = []     # (date, cust, name)
    src_by_grade = {}
    pruned = []       # 엑셀에서 사라져 제거한 기존 행 [(name, date)]
    page_orig = []    # 원본 각 페이지 데이터 높이 추정 → 안전 용량 산정
    for slide, tsh in dinfo:
        t = tsh.table
        trs = t._tbl.findall(qn("a:tr"))
        ph = 0
        for ri in range(2, len(t.rows)):
            cells = t.rows[ri].cells
            name = cells[0].text.strip()
            if not name:
                continue
            d = _norm_date(cells[1].text)
            cust = cells[3].text.strip()
            # 엑셀 등급확정 목록에 없는 기존 행은 제거(등급확정 해제/삭제된 건)
            if confirmed is not None and not any(
                    _same_incident(d, cust, name, cr) for cr in confirmed):
                pruned.append((name, cells[1].text.strip()))
                continue
            g = _grade_of(name)
            entries.append([g, d, copy.deepcopy(trs[ri])])
            existing.append((d, cust, name))
            src_by_grade.setdefault(g, trs[ri])
            ph += _tr_est_height(trs[ri])
        page_orig.append(ph)
    if pruned:
        log(f"  [슬라이드6~9] 엑셀에서 사라진(등급확정 해제) {len(pruned)}건을 상세목록에서 제거:")
        for nm, dt_ in pruned:
            log(f"       − {nm} ({dt_})")
    added = 0
    for it in new_items:
        d = _norm_date(it["date"])
        if any(_same_incident(d, it["cust"], it["name"],
                              {"date": ex[0], "cust": ex[1], "name": ex[2]}) for ex in existing):
            continue
        g = _grade_of(it["name"])
        src = src_by_grade.get(g)
        if src is None:
            src = src_by_grade.get(3)
        if src is None:
            src = next(iter(src_by_grade.values()))
        tr = copy.deepcopy(src)
        vals = [it["name"], it["date"], it["mins"], it["cust"], it["content"],
                it["cause"], it["div"], it["charge"], it["team"]]
        tcs = tr.findall(qn("a:tc"))
        for ci, v in enumerate(vals):
            if ci < len(tcs):
                _set_tc_text(tcs[ci], v)
        entries.append([g, d, tr])
        added += 1
    # 정렬: 1등급 → 2등급 → '3등급 이하'(3·4·…를 한 그룹) 순, 각 그룹 내에서는 발생일 오름차순.
    #   등급 부분을 min(등급,3)으로 묶어 3등급 이하는 등급 구분 없이 날짜순 정렬(예: 4등급 7/2 < 3등급 7/10).
    entries.sort(key=lambda e: (min(e[0], 3), e[1] if e[1] else (99, 99)))
    # 페이지 그룹핑: 안전 용량(원본 최대 페이지 높이) 초과 시 다음 페이지로
    budget = max(page_orig) if page_orig else 0
    groups, cur, used = [], [], 0
    for e in entries:
        h = _tr_est_height(e[2])
        if cur and budget and used + h > budget:
            groups.append(cur)
            cur, used = [], 0
        cur.append(e)
        used += h
    if cur:
        groups.append(cur)
    # 페이지 수 확보: 부족하면 마지막 상세 슬라이드 복제
    last_slide = dinfo[-1][0]
    while len(dinfo) < len(groups):
        dinfo.append(_duplicate_detail_slide(prs, last_slide))
    # 각 상세표에 그룹 배치
    for i, (slide, tsh) in enumerate(dinfo):
        _clear_data_rows(tsh.table)
        if i < len(groups):
            tbl = tsh.table._tbl
            for e in groups[i]:
                tbl.append(e[2])
    if len(groups) > len(page_orig):
        log(f"  [슬라이드] 상세 목록 용량 초과 → 상세 페이지 {len(groups) - len(page_orig)}장 추가")
    return added


def set_subtitle(sh, base_date):
    """0·1p 부제목: 본문('…발생 추이')은 20pt·볼드, 괄호('(M/D일 기준…)')는 16pt·일반으로 재구성.
    날짜/연도는 기준일로 치환하고, 원본 run의 폰트(설정된 영문/한글 서체)를 유지."""
    full = sub_date(sh.text_frame.text, base_date)
    idx = full.find("(")
    if idx < 0:
        set_shape_text(sh, full)
        return
    main_txt, paren_txt = full[:idx], full[idx:]
    p = sh.text_frame.paragraphs[0]
    pEl = p._p
    main_rpr = paren_rpr = None
    for r in p.runs:
        rpr = r._r.find(qn("a:rPr"))
        if rpr is None:
            continue
        if rpr.get("sz") == "1600" and paren_rpr is None:
            paren_rpr = copy.deepcopy(rpr)          # 괄호용(16pt·일반)
        elif rpr.get("sz") != "1600" and main_rpr is None:
            main_rpr = copy.deepcopy(rpr)           # 본문용(20pt·볼드)
    if main_rpr is None and p.runs:
        main_rpr = copy.deepcopy(p.runs[0]._r.find(qn("a:rPr")))
    if paren_rpr is None:
        paren_rpr = copy.deepcopy(main_rpr)
        if paren_rpr is not None:
            paren_rpr.set("sz", "1600")
            paren_rpr.set("b", "0")
    for r in list(p.runs):                          # 기존 run 제거
        r._r.getparent().remove(r._r)
    end = pEl.find(qn("a:endParaRPr"))

    def _mk(txt, rpr):
        r = pEl.makeelement(qn("a:r"), {})
        if rpr is not None:
            r.append(copy.deepcopy(rpr))
        t = r.makeelement(qn("a:t"), {})
        t.text = txt
        r.append(t)
        if end is not None:
            end.addprevious(r)
        else:
            pEl.append(r)

    _mk(main_txt, main_rpr)
    _mk(paren_txt, paren_rpr)


def set_title_with_count(sh, count_text, count_sz="1400"):
    """제목 갱신: 본문(부제목)은 원 서식 유지, 뒤 '(N건)'은 count_sz(기본 14pt)로.
    4·5페이지 부제목(신규 등급 확정 / 신규 발생 / 등급 협의 중)용."""
    p = sh.text_frame.paragraphs[0]
    pEl = p._p
    cur = "".join(r.text for r in p.runs)
    idx = cur.find("(")
    prefix = cur[:idx] if idx >= 0 else cur
    main_rpr = cnt_rpr = None
    for r in p.runs:
        rpr = r._r.find(qn("a:rPr"))
        if rpr is None:
            continue
        if rpr.get("sz") == count_sz and cnt_rpr is None:
            cnt_rpr = copy.deepcopy(rpr)          # 괄호용(14pt)
        elif rpr.get("sz") != count_sz and main_rpr is None:
            main_rpr = copy.deepcopy(rpr)          # 본문용(20pt)
    if main_rpr is None and p.runs:
        main_rpr = copy.deepcopy(p.runs[0]._r.find(qn("a:rPr")))
    if cnt_rpr is None:
        cnt_rpr = copy.deepcopy(main_rpr)
        if cnt_rpr is not None:
            cnt_rpr.set("sz", count_sz)
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    end = pEl.find(qn("a:endParaRPr"))

    def _mk(txt, rpr):
        r = pEl.makeelement(qn("a:r"), {})
        if rpr is not None:
            r.append(copy.deepcopy(rpr))
        t = r.makeelement(qn("a:t"), {})
        t.text = txt
        r.append(t)
        if end is not None:
            end.addprevious(r)
        else:
            pEl.append(r)

    _mk(prefix, main_rpr)
    _mk(count_text, cnt_rpr)


def set_detail_title(sh, count_text):
    """'장애 목록 상세' 제목 갱신.
    - 부제목('장애 목록 상세')은 원 서식(20pt, 설정된 한글 서체) 유지
    - 뒤 등급건수 텍스트는 원본 '건수 run'의 서식(14pt, 숫자=영문 서체 / 한글=한글 서체)을
      그대로 재사용 → 새 run을 만들며 한글 폰트를 굴림으로 폴백시키던 문제 방지."""
    p = sh.text_frame.paragraphs[0]
    runs = p.runs
    if not runs:
        p.add_run().text = "장애 목록 상세 " + count_text
        return
    runs[0].text = "장애 목록 상세 "
    if len(runs) >= 2:
        runs[1].text = count_text                     # 원본 건수 run 서식(latin/ea) 유지
        for r in list(runs[2:]):
            r._r.getparent().remove(r._r)
    else:
        # 건수 run이 없으면 프리픽스 run 복제 후 14pt + 폰트(숫자=영문 서체 / 한글=한글 서체) 지정
        newr = copy.deepcopy(runs[0]._r)
        rpr = newr.find(qn("a:rPr"))
        if rpr is not None:
            rpr.set("sz", "1400")
            for tag, face in (("a:latin", appconfig.FONT_LATIN), ("a:ea", appconfig.FONT_EA)):
                el = rpr.find(qn(tag))
                if el is None:
                    el = rpr.makeelement(qn(tag), {})
                    rpr.append(el)
                el.set("typeface", face)
        for t in newr.findall(qn("a:t")):
            t.text = count_text
        runs[0]._r.addnext(newr)


def _delete_confirmed_rows(table, confirmed):
    """표의 데이터행 중 confirmed(이번주 등급확정) 목록과 동일한 장애 행을 삭제. 삭제 수 반환."""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    remove = []
    for ri in range(2, len(table.rows)):
        cells = table.rows[ri].cells
        name = cells[0].text.strip()
        if not name:
            continue
        rdate = _norm_date(cells[1].text)
        rcust = cells[3].text.strip()
        if any(_same_incident(rdate, rcust, name, c) for c in confirmed):
            remove.append(ri)
    for ri in sorted(remove, reverse=True):
        tbl.remove(trs[ri])
    return len(remove)


def _set_count_title(prs, keyword, table):
    """제목의 '(N건)'을 표의 데이터행 수로 갱신(괄호 건수는 14pt)."""
    cnt = max(0, len(table.rows) - 2)
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and keyword in sh.text_frame.text:
                set_title_with_count(sh, f"({cnt}건)")


def _ln_width(ln):
    try:
        return int(ln.get("w") or "0")
    except Exception:
        return 0


def _set_line(ln, width, light):
    """선(a:lnX) 을 지정 두께/색(tx1, light=True면 lumMod/lumOff 50000 회색)으로 설정."""
    ln.set("w", str(width))
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:pattFill"):
        e = ln.find(qn(tag))
        if e is not None:
            ln.remove(e)
    sf = ln.makeelement(qn("a:solidFill"), {})
    sc = sf.makeelement(qn("a:schemeClr"), {"val": "tx1"})
    if light:
        sc.append(sc.makeelement(qn("a:lumMod"), {"val": "50000"}))
        sc.append(sc.makeelement(qn("a:lumOff"), {"val": "50000"}))
    sf.append(sc)
    ln.insert(0, sf)


def normalize_table_borders(prs):
    """모든 표의 내부 행 구분선을 가늘고 연하게(6350, 50% 회색) 통일.
    마지막 행의 아래 테두리는 바깥 테두리이므로 두껍게(12700) 유지/복원."""
    THIN, THICK, LIMIT = 6350, 12700, 9525
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_table:
                continue
            trs = sh.table._tbl.findall(qn("a:tr"))
            for i, tr in enumerate(trs):
                is_last = (i == len(trs) - 1)
                for tc in tr.findall(qn("a:tc")):
                    tcPr = tc.find(qn("a:tcPr"))
                    if tcPr is None:
                        continue
                    lnB = tcPr.find(qn("a:lnB"))
                    if lnB is None:
                        continue
                    if is_last:
                        if _ln_width(lnB) < LIMIT:      # 삭제 등으로 얇아진 바깥 하단 복원
                            _set_line(lnB, THICK, light=False)
                    elif _ln_width(lnB) >= LIMIT:       # 내부 두꺼운 선 → 가늘고 연하게
                        _set_line(lnB, THIN, light=True)


def fill_new_confirmed(out_prs, prev_prs, ws, base_date, do_summary=True):
    """슬라이드 4 '신규 등급 확정 장애' 표를 채운다. 작성된 항목 리스트 반환.
    do_summary=True 이면 엑셀 '장애 내용'을 규칙 기반으로 요약해 표에 넣는다."""
    C = DATA_COLS
    detail = []
    for t in _find_tables_by_title(prev_prs, "장애 목록 상세"):
        detail += _table_records(t)
    ts = _find_tables_by_title(prev_prs, "신규 발생 장애")
    sinbal = _table_records(ts[0]) if ts else []
    th = _find_tables_by_title(prev_prs, "등급 협의 중 장애")
    hyeobui = _table_records(th[0]) if th else []

    confirmed = []   # 이번주 등급확정 전체 (신규발생/협의중 삭제 매칭용)
    new_items = []
    for ri in range(2, ws.max_row + 1):
        if ws.cell(ri, C["year"]).value != base_date.year:
            continue
        start = ws.cell(ri, C["start"]).value
        if not isinstance(start, dt.datetime):
            continue
        d = (start.month, start.day)
        cust = str(ws.cell(ri, C["cust"]).value or "").strip()
        title = str(ws.cell(ri, C["title"]).value or "").strip()
        grade = str(ws.cell(ri, C["grade"]).value or "").strip()
        confirmed.append({"date": d, "cust": cust, "name": f"[{grade}] {title}"})

        def lookup(recs):
            for r in recs:
                if _same_incident(d, cust, title, r):
                    return r["content"]
            return None

        # 지난주 상세목록(6p~)에 이미 있으면 신규 아님
        if any(_same_incident(d, cust, title, r) for r in detail):
            continue
        # 장애내용 결정:
        #   원문 길이로 '지난주 PPT 텍스트 vs 엑셀 N열' 중 더 정보량 많은 쪽을 고르되,
        #   엑셀 쪽을 쓸 때는 규칙 기반 요약본을 넣는다(do_summary=True).
        prev_content, prev_src = lookup(sinbal), "지난주 신규발생"
        if prev_content is None:
            prev_content, prev_src = lookup(hyeobui), "지난주 협의중"
        excel_raw = str(ws.cell(ri, C["content"]).value or "").strip()
        excel_sum, low_conf = (summarize_content(excel_raw) if do_summary
                               else (excel_raw, False))
        use_excel = prev_content is None or len((prev_content or "").strip()) < len(excel_raw)
        if use_excel:
            content = excel_sum
            src = "엑셀(요약)" if (do_summary and len(excel_raw) >= 80) else "엑셀 데이터"
            if do_summary and low_conf and len(excel_raw) >= 80:
                src += " ⚠️저신뢰(수동검토)"
        else:
            content, src = prev_content, f"{prev_src}(유지)"
        new_items.append({
            "name": f"[{grade}] {title}", "date": f"{d[0]}/{d[1]}",
            "mins": ws.cell(ri, C["mins"]).value, "cust": cust, "content": content,
            "cause": str(ws.cell(ri, C["cause"]).value or "").strip(),
            "div": str(ws.cell(ri, C["div"]).value or "").strip(),
            "charge": str(ws.cell(ri, C["charge"]).value or "").strip(),
            "team": str(ws.cell(ri, C["team"]).value or "").strip(),
            "src": src,
        })

    tabs = _find_tables_by_title(out_prs, "신규 등급 확정 장애")
    if not tabs:
        log("  [경고] '신규 등급 확정 장애' 표를 찾지 못했습니다.")
        return new_items, confirmed
    table = tabs[0]
    _set_table_data_rows(table, max(len(new_items), 1))
    if new_items:
        for i, it in enumerate(new_items):
            cells = table.rows[2 + i].cells
            vals = [it["name"], it["date"], it["mins"], it["cust"], it["content"],
                    it["cause"], it["div"], it["charge"], it["team"]]
            for ci, v in enumerate(vals):
                _set_cell_text(cells[ci], v)
    else:
        for ci in range(len(table.rows[2].cells)):
            _set_cell_text(table.rows[2].cells[ci], "")

    # 제목의 '(N건)' 갱신
    for slide in out_prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and "신규 등급 확정 장애" in sh.text_frame.text:
                set_title_with_count(sh, f"({len(new_items)}건)")   # 괄호 건수 14pt

    # 이번주 등급확정으로 넘어간 건은 '신규 발생 장애'/'등급 협의 중 장애' 표에서 삭제
    for kw in ("신규 발생 장애", "등급 협의 중 장애"):
        tabs = _find_tables_by_title(out_prs, kw)
        if tabs:
            removed = _delete_confirmed_rows(tabs[0], confirmed)
            _set_count_title(out_prs, kw, tabs[0])
            if removed:
                log(f"  [{kw}] 등급확정 이동 {removed}건 삭제 (잔여 {len(tabs[0].rows) - 2}건)")

    # 신규확정 표가 커진 만큼 아래 '신규 발생 장애' 섹션을 내려 겹침 방지
    prev_conf = _find_tables_by_title(prev_prs, "신규 등급 확정 장애")
    if prev_conf:
        delta = _est_table_height(table) - _est_table_height(prev_conf[0])
        target = None
        for slide in out_prs.slides:
            if any(sh.has_table and sh.table._tbl is table._tbl for sh in slide.shapes):
                target = slide
                break
        if target and delta > 5000:
            title = next((sh for sh in target.shapes
                          if sh.has_text_frame and "신규 발생 장애" in sh.text_frame.text), None)
            if title is not None:
                base = title.top
                for sh in target.shapes:
                    if sh.top is not None and sh.top >= base - 1000:
                        sh.top = sh.top + delta
                log(f"  [슬라이드4] '신규 발생 장애' 섹션 {delta/914400:.2f}in 아래로 이동(겹침 방지)")
    return new_items, confirmed


# ─────────────────────────────────────────────────────────────────────────────
# 메인
# ─────────────────────────────────────────────────────────────────────────────
def parse_args():
    here = os.path.dirname(os.path.abspath(__file__))
    p = argparse.ArgumentParser(description="장애 주간보고 PPT 생성")
    p.add_argument("--xlsx", default=None,
                   help="업데이트된 엑셀 경로(미지정 시 output에서 자동 탐색)")
    p.add_argument("--template", default=None,
                   help="PPT 템플릿 경로(미지정 시 --prev 를 템플릿으로 사용)")
    p.add_argument("--prev", required=True,
                   help="지난주 장애공유 PPT 경로 (템플릿 겸 비교용, 필수)")
    p.add_argument("--date", default=None,
                   help="기준일 (YYYY-MM-DD). 기본: 엑셀 기준일(요약!B2), 없으면 오늘")
    p.add_argument("--date-source", choices=("auto", "excel", "input"), default="auto",
                   help="기준일 우선순위: auto(엑셀 우선)·excel(무조건 엑셀)·input(무조건 --date)")
    p.add_argument("--out", default=None, help="결과 PPT 경로(미지정 시 지난주 PPT 폴더에 생성)")
    p.add_argument("--outdir", default=None, help="결과 폴더(미지정 시 지난주 PPT 폴더)")
    p.add_argument("--raw-content", action="store_true",
                   help="엑셀 '장애 내용'을 요약하지 않고 원문 그대로 넣음(기본: 규칙 기반 요약)")
    return p.parse_args()


def resolve_xlsx(args, base_date, here):
    if args.xlsx:
        return os.path.abspath(args.xlsx)
    # 인플레이스 갱신된 원본 엑셀(파일명은 config)
    cand = os.path.join(here, "data", appconfig.EXCEL_NAME)
    if os.path.exists(cand):
        return cand
    # (구버전 호환) output 폴더의 최신 _updated_ 파일
    files = glob.glob(os.path.join(here, "output", "*_updated_*.xlsx"))
    return max(files, key=os.path.getmtime) if files else None


def read_basedate(xlsx):
    """엑셀 '요약' 시트의 기준일(B2, 없으면 B1)을 date-only datetime으로 반환. 없으면 None.
    (런처가 in-process로 호출해 화면 날짜 자동 채움/불일치 확인에 사용)"""
    wb = None
    try:
        wb = openpyxl.load_workbook(xlsx, data_only=True, read_only=True)
        ws = wb["요약"]
        for addr in ("B2", "B1"):
            v = ws[addr].value
            if isinstance(v, dt.datetime):
                return dt.datetime(v.year, v.month, v.day)
    except Exception:
        return None
    finally:
        if wb is not None:
            try:
                wb.close()
            except Exception:
                pass
    return None


def main():
    args = parse_args()
    here = os.path.dirname(os.path.abspath(__file__))

    input_date = dt.datetime.strptime(args.date, "%Y-%m-%d") if args.date else None

    xlsx = resolve_xlsx(args, input_date, here)
    if not xlsx or not os.path.exists(xlsx):
        log("[오류] 업데이트된 엑셀을 찾을 수 없습니다.")
        log("       먼저 update_excel.py 를 실행하거나 --xlsx 로 경로를 지정하세요.")
        sys.exit(1)

    # 리포트 기준일 결정: 원칙은 '엑셀 기준일'(차트·건수가 엑셀에서 나오므로 제목도 데이터와 일치).
    #   --date-source: auto(엑셀 우선) / excel(무조건 엑셀) / input(무조건 --date)
    excel_date = read_basedate(xlsx)
    today = dt.date.today()
    if args.date_source == "input" and input_date:
        base_date = input_date
    elif args.date_source == "excel" and excel_date:
        base_date = excel_date
    else:  # auto
        base_date = excel_date or input_date or dt.datetime(today.year, today.month, today.day)
    by, bm = base_date.year, base_date.month

    if excel_date and input_date and excel_date.date() != input_date.date():
        if base_date.date() == excel_date.date():
            log(f"  [알림] 입력 기준일({input_date:%Y-%m-%d}) != 엑셀 기준일({excel_date:%Y-%m-%d})"
                f" -> 데이터와 일치하도록 엑셀 기준일 사용. (입력값 강제: --date-source input)")
        else:
            log(f"  [경고] 입력 기준일({input_date:%Y-%m-%d}) 사용 — 엑셀 데이터는"
                f" {excel_date:%Y-%m-%d} 기준이라 제목과 실제 데이터 날짜가 다릅니다.")
    elif not excel_date:
        log("  [알림] 엑셀 기준일(요약!B2)을 읽지 못해 입력/오늘 날짜를 사용합니다.")
    if not os.path.exists(args.prev):
        log(f"[오류] 지난주 PPT를 찾을 수 없습니다: {args.prev}")
        sys.exit(1)
    template = args.template or args.prev      # 템플릿 미지정 시 지난주 PPT를 템플릿으로 사용
    if not os.path.exists(template):
        log(f"[오류] 템플릿을 찾을 수 없습니다: {template}")
        sys.exit(1)

    if args.out:
        out = os.path.abspath(args.out)
    else:
        outdir = args.outdir or os.path.dirname(os.path.abspath(args.prev))
        os.makedirs(outdir, exist_ok=True)
        out = os.path.join(outdir, f"{base_date:%Y%m%d}_장애공유.pptx")

    log("=" * 70)
    log("장애 주간보고 - PPT 생성")
    log("=" * 70)
    log(f"  엑셀     : {xlsx}")
    log(f"  템플릿   : {template}" + ("  (지난주 PPT 사용)" if not args.template else ""))
    log(f"  기준일   : {base_date:%Y-%m-%d}  (최근 3개년 {by-2}~{by})")
    log(f"  결과 PPT : {out}")
    log("-" * 70)

    wb = openpyxl.load_workbook(xlsx, data_only=True)
    ws_sum = wb["요약"]

    # 최근 3개년 월 범위
    months = month_range(by - 2, 1, by, bm)
    categories = [f"{m}월" for (_, m) in months]
    log(f"  차트 구간: {by-2}-01 ~ {by}-{bm:02d}  ({len(months)}개월)")

    # 사업부별 시리즈 읽기
    def series_for(cfg_key):
        cfg = SRC[cfg_key]
        ws = wb[cfg["sheet"]]
        total = read_series(ws, cfg["total_row"], cfg["total_start_col"], months)
        divs = {name: read_series(ws, row, cfg["div_start_col"], months)
                for name, row in cfg["divisions"].items()}
        return total, divs

    total0, divs0 = series_for("사업부별")            # 전사 / 사업부별
    total1, divs1 = series_for("사업부별_계열사별")    # 계열사 전사 / 계열사 사업부별

    prs = Presentation(template)

    # ── 슬라이드 0: 전사 ──
    s0 = prs.slides[0]
    for sh in s0.shapes:
        if sh.has_chart:
            n_tpl = len(list(sh.chart.plots[0].categories))
            replace_chart(sh, categories, total0)
            align_year_bands(s0, n_tpl, len(months), by)
            log(f"  [슬라이드0] 전사 추이 차트 갱신 (합계 {sum(total0)}건) + 연도 음영/라벨 1월 정렬")
    p19 = find_shape(s0, 19)
    if p19:
        paras = p19.text_frame.paragraphs
        set_paragraph_text(paras[0], group_commas(cell(ws_sum, "C22")))
        if len(paras) > 1:
            set_paragraph_text(paras[1], group_commas(cell(ws_sum, "C23")))
        log("  [슬라이드0] 요약문구 갱신 (요약!C22/C23)")
    s0_4 = find_shape(s0, 4)
    if s0_4:
        set_shape_text(s0_4, polish_delta(cell(ws_sum, "C24")))
    sh = find_shape(s0, 2)          # 부제목: 본문 20pt·볼드 / 괄호 16pt·일반
    if sh:
        set_subtitle(sh, base_date)
    sh = find_shape(s0, 23)         # 우상단 '(M/D 기준)' 작은 라벨
    if sh:
        set_shape_text(sh, sub_date(sh.text_frame.text, base_date))

    # ── 슬라이드 1: 계열사 전사 ──
    s1 = prs.slides[1]
    for sh in s1.shapes:
        if sh.has_chart:
            n_tpl = len(list(sh.chart.plots[0].categories))
            replace_chart(sh, categories, total1)
            align_year_bands(s1, n_tpl, len(months), by)
            log(f"  [슬라이드1] 계열사 전사 추이 차트 갱신 (합계 {sum(total1)}건) + 연도 음영/라벨 1월 정렬")
    p19 = find_shape(s1, 19)
    if p19:
        paras = p19.text_frame.paragraphs
        set_paragraph_text(paras[0], group_commas(cell(ws_sum, "C135")))
        if len(paras) > 1:
            set_paragraph_text(paras[1], group_commas(cell(ws_sum, "C136")))
        log("  [슬라이드1] 요약문구 갱신 (요약!C135/C136)")
    s1_4 = find_shape(s1, 4)
    if s1_4:
        set_shape_text(s1_4, polish_delta(cell(ws_sum, "C137")))
    sh = find_shape(s1, 2)          # 부제목: 본문 20pt·볼드 / 괄호 16pt·일반
    if sh:
        set_subtitle(sh, base_date)
    sh = find_shape(s1, 23)         # 우상단 '(M/D 기준)' 작은 라벨
    if sh:
        set_shape_text(sh, sub_date(sh.text_frame.text, base_date))

    # ── 슬라이드 2,3: 사업부별 차트 (제목=사업부명으로 매핑) ──
    def update_div_slide(slide, divs, label):
        cnt = 0
        for sh in slide.shapes:
            if not sh.has_chart:
                continue
            ch = sh.chart
            title = ch.chart_title.text_frame.text if ch.has_title else None
            if title in divs:
                replace_chart(sh, categories, divs[title])
                cnt += 1
        log(f"  [{label}] 사업부 차트 {cnt}개 갱신")

    update_div_slide(prs.slides[2], divs0, "슬라이드2")
    n2 = align_year_labels_grid(prs.slides[2], len(months), by)
    log(f"  [슬라이드2] 연도 라벨 1월 정렬 ({n2}개)")
    update_div_slide(prs.slides[3], divs1, "슬라이드3")
    n3 = align_year_labels_grid(prs.slides[3], len(months), by)
    log(f"  [슬라이드3] 연도 라벨 1월 정렬 ({n3}개)")

    # ── 슬라이드 6~9: 제목 등급 건수 (요약!C22에서 추출) ──
    c22 = cell(ws_sum, "C22")
    m_tot = re.search(r"총\s*([\d,]+)\s*건", c22)
    m_g2 = re.search(r"2등급\s*([\d,]+)\s*건", c22)
    m_g3 = re.search(r"3등급\s*이하\s*([\d,]+)\s*건", c22)
    detail_count_text = None    # 상세목록 제목 건수 (rebuild 이후 모든 상세 슬라이드에 적용)
    if m_tot and m_g2 and m_g3:
        tot = m_tot.group(1); g2 = m_g2.group(1); g3 = m_g3.group(1)
        detail_count_text = group_commas(f"(2등급 {g2}건, 3등급 이하 {g3}건 총 {tot}건)")
    else:
        log("  [경고] 요약!C22에서 등급 건수를 추출하지 못해 제목을 갱신하지 않음")

    # ── 슬라이드 4: 신규 등급 확정 장애 (지난주 PPT 비교) ──
    prev_prs = Presentation(args.prev)
    new_items, confirmed = fill_new_confirmed(prs, prev_prs, wb["데이터"], base_date,
                                              do_summary=not args.raw_content)
    log(f"  [슬라이드4] 신규 등급 확정 장애 {len(new_items)}건 작성 (지난주: {os.path.basename(args.prev)})"
        + ("  [내용요약: ON]" if not args.raw_content else "  [내용요약: OFF(원문)]"))
    for it in new_items:
        log(f"     - {it['name']}  (발생 {it['date']}, 내용출처: {it['src']})")
    low = [it for it in new_items if "저신뢰" in it["src"]]
    if low:
        log(f"  [내용요약] ⚠️ 규칙 신뢰도 낮음 {len(low)}건 — PPT 생성 후 아래 항목의 '장애내용' 칸을 눈으로 확인하세요:")
        for it in low:
            log(f"       · {it['name']} (발생 {it['date']})")

    # ── 슬라이드 6~9: '장애 목록 상세'에 신규확정 누락건 추가 + 엑셀 삭제건 제거 + (등급→발생일) 재정렬 ──
    added = rebuild_detail_list(prs, new_items, confirmed)
    log(f"  [슬라이드6~9] 장애 목록 상세 재정렬 완료 (신규확정 {added}건 추가)")

    # 상세목록 제목 등급건수를 '모든' 상세 슬라이드에 적용(페이지 추가/삭제 후 실행 → 인덱스 하드코딩 안 함)
    if detail_count_text:
        n = 0
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.has_text_frame and "장애 목록 상세" in sh.text_frame.text:
                    set_detail_title(sh, detail_count_text); n += 1
        log(f"  [슬라이드6~] 제목 등급건수 갱신(14pt) {n}개: 장애 목록 상세 {detail_count_text}")

    # 모든 표의 내부 행 구분선을 가늘고 연하게 통일 (행 추가/삭제로 생긴 두꺼운 선 정리)
    normalize_table_borders(prs)
    log("  [전체] 표 내부 행 구분선 정규화 완료")

    prs.save(out)
    log("-" * 70)
    log(f"  저장 완료: {out}")
    log("=" * 70)
    log("완료.")
    log("=" * 70)


if __name__ == "__main__":
    main()
